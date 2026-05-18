import os
import pandas as pd
from pathlib import Path
from typing import List, Dict, Any
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup
import markdown
from src.models.document import Document, DocumentMetadata
from src.utils.logger import logger
from src.processors.pdf_parser import PDFParser

class MultiDocParser:
    def __init__(self, use_ocr: bool = False):
        self.pdf_parser = PDFParser(use_ocr=use_ocr)

    def parse(self, file_path: str) -> Document:
        path = Path(file_path)
        suffix = path.suffix.lower()
        
        if suffix == ".pdf":
            return self.pdf_parser.parse(file_path)
        elif suffix in [".docx", ".doc"]:
            return self._parse_docx(file_path)
        elif suffix in [".xlsx", ".xls", ".csv"]:
            return self._parse_excel(file_path)
        elif suffix in [".pptx", ".ppt"]:
            return self._parse_pptx(file_path)
        elif suffix == ".md":
            return self._parse_markdown(file_path)
        elif suffix in [".html", ".htm"]:
            return self._parse_html(file_path)
        elif suffix == ".txt":
            return self._parse_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {suffix}")

    def _parse_docx(self, file_path: str) -> Document:
        try:
            from docx.table import Table
            from docx.text.paragraph import Paragraph
            
            doc = DocxDocument(file_path)
            full_text = []
            
            # Iterate over elements in order
            for element in doc.element.body.iterchildren():
                if element.tag.endswith('p'):
                    para = Paragraph(element, doc)
                    if para.text.strip():
                        full_text.append(para.text)
                
                elif element.tag.endswith('tbl'):
                    table = Table(element, doc)
                    table_rows = []
                    for row in table.rows:
                        # Use a simple separator for cells
                        row_cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                        table_rows.append(" | ".join(row_cells))
                    
                    if table_rows:
                        full_text.append("\n[Table]\n" + "\n".join(table_rows) + "\n")
            
            combined_text = "\n\n".join(full_text)
            return self._create_document(file_path, combined_text, "docx")
        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {e}")
            return self._failed_document(file_path, e)

    def _parse_excel(self, file_path: str) -> Document:
        try:
            suffix = Path(file_path).suffix.lower()
            if suffix == ".csv":
                df = pd.read_csv(file_path)
                combined_text = df.to_string()
            else:
                all_sheets = pd.read_excel(file_path, sheet_name=None)
                parts = []
                for i, (sheet_name, df) in enumerate(all_sheets.items()):
                    parts.append(f"[Page {i+1}] Sheet: {sheet_name}\n{df.to_string()}")
                combined_text = "\n\n".join(parts)
            
            return self._create_document(file_path, combined_text, "excel")
        except Exception as e:
            logger.error(f"Error parsing Excel/CSV {file_path}: {e}")
            return self._failed_document(file_path, e)

    def _parse_pptx(self, file_path: str) -> Document:
        try:
            prs = Presentation(file_path)
            full_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = [f"[Slide {i+1}]"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                full_text.append("\n".join(slide_text))
            
            combined_text = "\n\n".join(full_text)
            return self._create_document(file_path, combined_text, "pptx")
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            return self._failed_document(file_path, e)

    def _parse_markdown(self, file_path: str) -> Document:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            # Simple markdown parsing, could be enhanced with metadata extraction
            return self._create_document(file_path, text, "markdown")
        except Exception as e:
            logger.error(f"Error parsing MD {file_path}: {e}")
            return self._failed_document(file_path, e)

    def _parse_html(self, file_path: str) -> Document:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f, "html.parser")
            
            # Remove scripts and styles
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text(separator="\n\n")
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            combined_text = "\n".join(chunk for chunk in chunks if chunk)
            
            return self._create_document(file_path, combined_text, "html")
        except Exception as e:
            logger.error(f"Error parsing HTML {file_path}: {e}")
            return self._failed_document(file_path, e)

    def _parse_txt(self, file_path: str) -> Document:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            return self._create_document(file_path, text, "txt")
        except Exception as e:
            logger.error(f"Error parsing TXT {file_path}: {e}")
            return self._failed_document(file_path, e)

    def _create_document(self, file_path: str, content: str, file_type: str) -> Document:
        path = Path(file_path)
        metadata = DocumentMetadata(
            title=path.stem,
            source=str(path),
            file_type=file_type
        )
        return Document(
            filename=path.name,
            content=content,
            metadata=metadata,
            status="processed"
        )

    def _failed_document(self, file_path: str, e: Exception) -> Document:
        return Document(
            filename=Path(file_path).name,
            status="failed",
            error_message=str(e)
        )
